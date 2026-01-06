"""
PDF to PowerPoint Web Application - Backend (LIGHTWEIGHT VERSION)
FastAPI server with Gemini API integration for PDF analysis and PPTX generation
Optimized for low-memory environments (512MB RAM)
"""
import os
import shutil
import json
import uuid
import base64
import asyncio
import gc
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
import fitz  # PyMuPDF

# Gemini API
import google.generativeai as genai

# Import conversion modules
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

app = FastAPI(
    title="PDF to PowerPoint Converter",
    description="Convert PDF files to editable PowerPoint presentations using AI",
    version="1.0.0"
)

# CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Storage for job status
jobs = {}

# Directories
BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "output"
TEMP_DIR = BASE_DIR / "temp_processing"

# Create directories
for d in [UPLOAD_DIR, OUTPUT_DIR, TEMP_DIR]:
    d.mkdir(exist_ok=True)


class JobStatus:
    PENDING = "pending"
    PROCESSING = "processing"
    ANALYZING = "analyzing"
    GENERATING = "generating"
    COMPLETED = "completed"
    ERROR = "error"


# Gemini API Setup
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

# Image Analysis Prompt
IMAGE_ANALYSIS_PROMPT = """
あなたはPDFページ画像からテキスト情報を抽出するアナリストです。以下の画像を分析し、指定されたJSONフォーマットで出力してください。

## 出力フォーマット

```json
{
  "replace_all": true,
  "blocks": [
    {
      "text": "テキスト内容",
      "bbox_1000": [x, y, width, height],
      "font_family": "フォント名",
      "is_bold": true/false,
      "font_size_pt": 数値,
      "colors": [
        {"range": [開始文字位置, 終了文字位置], "rgb": [R, G, B]}
      ]
    }
  ]
}
```

## 詳細ルール

### 座標 (bbox_1000)
- 画像を1000x1000の座標系として扱う
- [x, y, width, height] の形式
- x: 左端からの距離 (0-1000)
- y: 上端からの距離 (0-1000)
- width, height: テキストボックスの幅と高さ

### フォントファミリー (font_family)
以下の8種類から選択：

**日本語フォント:**
- `Noto Sans JP` - ゴシック体
- `Noto Serif JP` - 明朝体
- `Yomogi` - 手書き風
- `Kosugi Maru` - 丸文字

**英語フォント:**
- `Roboto` - サンセリフ（標準）
- `Merriweather` - セリフ
- `Roboto Mono` - 等幅
- `Montserrat` - 太め見出し

### フォントサイズ (font_size_pt)
- PowerPointスライド（幅1376pt × 高さ768pt）基準
- font_size_pt = (テキスト高さ / 画像高さ) × 768

### テキストグループ化ルール
1. 縦方向（Y座標が異なる）→ 必ず別のblock
2. 横方向（同じ行）→ 距離が近ければ同一block
3. 色が異なる場合は`colors`配列で表現
4. 改行は使用禁止、別blockに分割

JSONのみを出力してください。説明文は不要です。
"""


def configure_gemini():
    """Configure Gemini API"""
    if GEMINI_API_KEY:
        genai.configure(api_key=GEMINI_API_KEY)
        return True
    return False


async def analyze_image_with_gemini(image_path: Path, page_num: int) -> dict:
    """Analyze a single page image with Gemini API"""
    if not configure_gemini():
        raise ValueError("GEMINI_API_KEY not set")
    
    # Read and encode image
    with open(image_path, "rb") as f:
        image_data = f.read()
    
    # Create Gemini model
    model = genai.GenerativeModel("gemini-3-flash-preview")
    
    # Create image part
    image_part = {
        "mime_type": "image/png",
        "data": image_data
    }
    
    # Generate content
    response = model.generate_content([IMAGE_ANALYSIS_PROMPT, image_part])
    
    # Parse response
    response_text = response.text.strip()
    
    # Extract JSON from response
    if "```json" in response_text:
        json_start = response_text.find("```json") + 7
        json_end = response_text.find("```", json_start)
        response_text = response_text[json_start:json_end].strip()
    elif "```" in response_text:
        json_start = response_text.find("```") + 3
        json_end = response_text.find("```", json_start)
        response_text = response_text[json_start:json_end].strip()
    
    try:
        page_data = json.loads(response_text)
        return page_data
    except json.JSONDecodeError as e:
        print(f"JSON parse error for page {page_num}: {e}")
        print(f"Response text: {response_text[:500]}")
        # Return placeholder on error
        return {
            "replace_all": True,
            "blocks": [{
                "text": f"[Page {page_num} - Parse error]",
                "bbox_1000": [50, 50, 900, 100],
                "font_family": "Roboto",
                "is_bold": True,
                "font_size_pt": 32,
                "colors": [{"range": [0, 30], "rgb": [30, 30, 30]}]
            }]
        }


@app.get("/")
async def root():
    api_status = "configured" if GEMINI_API_KEY else "not configured"
    return {"message": "PDF to PowerPoint Converter API", "status": "running", "gemini_api": api_status}


@app.post("/api/upload")
async def upload_pdf(file: UploadFile = File(...), mode: str = Form("precision")):
    """Upload a PDF file"""
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed")
    
    job_id = str(uuid.uuid4())
    job_dir = TEMP_DIR / job_id
    job_dir.mkdir(exist_ok=True)
    
    pdf_path = job_dir / "input.pdf"
    with open(pdf_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    # mode: 'precision' (standard) or 'safeguard' (with backup layer)
    jobs[job_id] = {
        "status": JobStatus.PENDING,
        "progress": 0,
        "message": "PDF uploaded successfully",
        "pdf_path": str(pdf_path),
        "job_dir": str(job_dir),
        "original_filename": file.filename,
        "total_pages": 0,
        "current_page": 0,
        "mode": mode  # 'precision' or 'safeguard'
    }
    
    return {"job_id": job_id, "message": "Upload successful", "mode": mode}


@app.post("/api/process/{job_id}")
async def start_processing(job_id: str, background_tasks: BackgroundTasks):
    """Start processing the uploaded PDF"""
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    if not GEMINI_API_KEY:
        raise HTTPException(status_code=500, detail="GEMINI_API_KEY not configured")
    
    background_tasks.add_task(process_pdf_with_gemini, job_id)
    jobs[job_id]["status"] = JobStatus.PROCESSING
    jobs[job_id]["message"] = "Processing started"
    
    return {"status": "processing", "message": "Processing started"}


async def process_pdf_with_gemini(job_id: str):
    """Background task to process PDF with Gemini API"""
    try:
        job = jobs[job_id]
        job_dir = Path(job["job_dir"])
        pdf_path = Path(job["pdf_path"])
        
        # Step 1: Convert PDF to images
        job["status"] = JobStatus.PROCESSING
        job["message"] = "Converting PDF to images..."
        
        doc = fitz.open(pdf_path)
        total_pages = len(doc)
        job["total_pages"] = total_pages
        
        pages_dir = job_dir / "pages"
        pages_dir.mkdir(exist_ok=True)
        
        page_width = doc[0].rect.width
        page_height = doc[0].rect.height
        job["page_width"] = page_width
        job["page_height"] = page_height
        
        # LIGHTWEIGHT: Use 1.5x scale instead of 2.0x (saves ~40% memory)
        for page_num in range(total_pages):
            page = doc[page_num]
            mat = fitz.Matrix(1.5, 1.5)  # Reduced from 2.0
            pix = page.get_pixmap(matrix=mat)
            img_path = pages_dir / f"page_{page_num + 1}.png"
            pix.save(str(img_path))
            
            # Clean up pixmap to free memory immediately
            del pix
            gc.collect()
            
            job["current_page"] = page_num + 1
            job["progress"] = int((page_num + 1) / total_pages * 20)
        
        doc.close()
        del doc
        gc.collect()
        
        # Step 2: Analyze images with Gemini
        job["status"] = JobStatus.ANALYZING
        job["message"] = "Analyzing page content with AI..."
        
        analysis_results = {}
        
        for page_num in range(1, total_pages + 1):
            job["message"] = f"Analyzing page {page_num}/{total_pages}..."
            job["progress"] = 20 + int(page_num / total_pages * 40)
            
            img_path = pages_dir / f"page_{page_num}.png"
            
            try:
                page_data = await analyze_image_with_gemini(img_path, page_num)
                analysis_results[f"page_{page_num}"] = page_data
            except Exception as e:
                print(f"Error analyzing page {page_num}: {e}")
                analysis_results[f"page_{page_num}"] = {
                    "replace_all": True,
                    "blocks": [{
                        "text": f"[Page {page_num} - Analysis error: {str(e)[:50]}]",
                        "bbox_1000": [50, 50, 900, 100],
                        "font_family": "Roboto",
                        "is_bold": True,
                        "font_size_pt": 24,
                        "colors": [{"range": [0, 50], "rgb": [200, 50, 50]}]
                    }]
                }
            # Clean up memory after each page analysis
            gc.collect()
            
            # Small delay to avoid rate limiting
            await asyncio.sleep(0.5)
        
        # Save analysis JSON
        json_path = job_dir / "image_analysis.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(analysis_results, f, ensure_ascii=False, indent=2)
        
        job["progress"] = 65
        
        # Step 3: Generate PowerPoint using standalone converter
        job["status"] = JobStatus.GENERATING
        job["message"] = "Generating PowerPoint..."
        
        output_filename = Path(job["original_filename"]).stem + ".pptx"
        output_path = OUTPUT_DIR / f"{job_id}_{output_filename}"
        
        # Select converter based on mode - LIGHTWEIGHT versions
        mode = job.get("mode", "precision")
        if mode == "safeguard":
            converter_script = BASE_DIR / "standalone_convert_v4_v43_light.py"
            print(f"Using Safeguard Mode converter (v43 LIGHT)")
        else:
            converter_script = BASE_DIR / "standalone_convert_v43_light.py"
            print(f"Using Precision Mode converter (v43 LIGHT)")
        
        log_path = job_dir / "conversion_log.txt"
        
        # Run converter as subprocess
        import subprocess
        cmd = [
            sys.executable,
            str(converter_script),
            "--pdf", str(pdf_path),
            "--output", str(output_path),
            "--json", str(json_path),
            "--log", str(log_path)
        ]
        
        print(f"Running converter: {' '.join(cmd)}")
        
        # Update progress periodically while subprocess runs
        process = subprocess.Popen(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            cwd=str(BASE_DIR)
        )
        
        # Wait for completion with progress updates
        while process.poll() is None:
            job["progress"] = min(95, job["progress"] + 1)
            await asyncio.sleep(2)
        
        returncode = process.returncode
        if returncode != 0:
            stderr = process.stderr.read().decode("utf-8", errors="replace")
            raise Exception(f"Converter failed with code {returncode}: {stderr[:500]}")
        
        job["status"] = JobStatus.COMPLETED
        job["progress"] = 100
        job["message"] = "Conversion completed!"
        job["output_path"] = str(output_path)
        job["output_filename"] = output_filename
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        jobs[job_id]["status"] = JobStatus.ERROR
        jobs[job_id]["message"] = str(e)
        jobs[job_id]["progress"] = 0


async def generate_pptx_from_analysis(
    pdf_path: Path, json_path: Path, pages_dir: Path, output_path: Path, 
    job: dict, page_width: float, page_height: float, total_pages: int
):
    """Generate PPTX from Gemini analysis"""
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from PIL import Image
    import cv2
    import numpy as np
    from collections import Counter
    
    def set_font_for_run(run, font_name):
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'))
            rPr.append(ea)
        ea.set('typeface', font_name)
        latin = rPr.find(qn('a:latin'))
        if latin is None:
            latin = rPr.makeelement(qn('a:latin'))
            rPr.insert(0, latin)
        latin.set('typeface', font_name)
    
    def get_background_color(image):
        w, h = image.size
        samples = [
            image.getpixel((10, 10)),
            image.getpixel((w-11, 10)),
            image.getpixel((10, h-11)),
            image.getpixel((w-11, h-11)),
        ]
        samples = [s[:3] if len(s) > 3 else s for s in samples]
        return Counter(samples).most_common(1)[0][0]
    
    # Load analysis JSON
    with open(json_path, "r", encoding="utf-8") as f:
        analysis = json.load(f)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Emu(int(page_width * 12700))
    prs.slide_height = Emu(int(page_height * 12700))
    
    for page_num in range(1, total_pages + 1):
        job["message"] = f"Generating slide {page_num}/{total_pages}..."
        job["progress"] = 65 + int(page_num / total_pages * 30)
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Get page image
        img_path = pages_dir / f"page_{page_num}.png"
        page_img = Image.open(img_path)
        img_width, img_height = page_img.size
        
        # Set background color
        bg_color = get_background_color(page_img)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
        
        # Get page data
        page_key = f"page_{page_num}"
        page_data = analysis.get(page_key, {"blocks": []})
        blocks = page_data.get("blocks", [])
        
        # Detect and add graphic elements (simplified version)
        # Use numpy for Japanese path support
        img_array = np.fromfile(str(img_path), dtype=np.uint8)
        cv_img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
        if cv_img is not None:
            # Detect non-text graphic areas
            bg_bgr = np.array([bg_color[2], bg_color[1], bg_color[0]])
            diff = cv2.absdiff(cv_img, np.full_like(cv_img, bg_bgr))
            diff_gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
            _, mask = cv2.threshold(diff_gray, 25, 255, cv2.THRESH_BINARY)
            
            kernel = np.ones((5, 5), np.uint8)
            mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel)
            
            contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Get text regions for exclusion
            text_regions = []
            for block in blocks:
                if "bbox_1000" in block:
                    nx, ny, nw, nh = block["bbox_1000"]
                    px = int(nx / 1000 * img_width)
                    py = int(ny / 1000 * img_height)
                    pw = int(nw / 1000 * img_width)
                    ph = int(nh / 1000 * img_height)
                    text_regions.append((px, py, px + pw, py + ph))
            
            # Add large graphic elements
            for contour in contours:
                area = cv2.contourArea(contour)
                if area < 1000:
                    continue
                
                x, y, w, h = cv2.boundingRect(contour)
                
                # Check text overlap
                is_text_area = False
                for tx1, ty1, tx2, ty2 in text_regions:
                    ox1 = max(x, tx1)
                    oy1 = max(y, ty1)
                    ox2 = min(x + w, tx2)
                    oy2 = min(y + h, ty2)
                    if ox2 > ox1 and oy2 > oy1:
                        overlap = (ox2 - ox1) * (oy2 - oy1)
                        if overlap / (w * h) > 0.3:
                            is_text_area = True
                            break
                
                if is_text_area:
                    continue
                
                # Extract and add element
                elem = cv_img[y:y+h, x:x+w].copy()
                elem_mask = mask[y:y+h, x:x+w].copy()
                elem_rgba = cv2.cvtColor(elem, cv2.COLOR_BGR2BGRA)
                elem_rgba[:, :, 3] = elem_mask
                
                elem_path = pages_dir / f"elem_{page_num}_{x}_{y}.png"
                cv2.imwrite(str(elem_path), elem_rgba)
                
                left_pt = x / img_width * page_width
                top_pt = y / img_height * page_height
                width_pt = w / img_width * page_width
                height_pt = h / img_height * page_height
                
                try:
                    slide.shapes.add_picture(str(elem_path), Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
                except:
                    pass
        
        # Add text blocks
        for block in blocks:
            text = block.get("text", "")
            if not text:
                continue
            
            # Get coordinates
            if "bbox_1000" in block:
                nx, ny, nw, nh = block["bbox_1000"]
                left_pt = nx / 1000 * page_width
                top_pt = ny / 1000 * page_height
                width_pt = nw / 1000 * page_width
                height_pt = nh / 1000 * page_height
            else:
                continue
            
            if width_pt < 10:
                width_pt = 50
            if height_pt < 10:
                height_pt = 30
            
            textbox = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
            tf = textbox.text_frame
            tf.word_wrap = False
            tf.margin_left = 0
            tf.margin_top = 0
            tf.margin_right = 0
            tf.margin_bottom = 0
            
            p = tf.paragraphs[0]
            
            font_size = block.get("font_size_pt", 12)
            font_family = block.get("font_family", "Roboto")
            is_bold = block.get("is_bold", False)
            colors = block.get("colors", [])
            
            if colors:
                p.text = ""
                covered = 0
                for color_info in colors:
                    start, end = color_info.get("range", [0, len(text)])
                    rgb = color_info.get("rgb", [0, 0, 0])
                    
                    run_text = text[start:end]
                    if run_text:
                        run = p.add_run()
                        run.text = run_text
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))
                        set_font_for_run(run, font_family)
                        run.font.bold = is_bold
                    covered = max(covered, end)
                
                if covered < len(text):
                    remaining = text[covered:]
                    run = p.add_run()
                    run.text = remaining
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    set_font_for_run(run, font_family)
                    run.font.bold = is_bold
            else:
                p.text = text
                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    set_font_for_run(run, font_family)
                    run.font.bold = is_bold
    
    prs.save(output_path)
    job["progress"] = 98


@app.get("/api/status/{job_id}")
async def get_status(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    return jobs[job_id]


@app.get("/api/download/{job_id}")
async def download_result(job_id: str):
    from urllib.parse import quote
    
    # Try to find job in memory first
    if job_id in jobs:
        job = jobs[job_id]
        if job["status"] != JobStatus.COMPLETED:
            raise HTTPException(status_code=400, detail="Processing not completed")
        
        output_path = Path(job["output_path"])
        output_filename = job["output_filename"]
    else:
        # Fallback: Search for file in output directory matching job_id
        possible_files = list(OUTPUT_DIR.glob(f"{job_id}*"))
        if not possible_files:
            raise HTTPException(status_code=404, detail="Job not found and no matching file in output directory")
        
        output_path = possible_files[0]
        output_filename = output_path.name
    
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="Output file not found")
    
    # Handle Japanese filenames for download display
    encoded_filename = quote(output_filename)
    
    return FileResponse(
        path=str(output_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{encoded_filename}"
        }
    )


@app.delete("/api/job/{job_id}")
async def cleanup_job(job_id: str):
    if job_id in jobs:
        job_dir = TEMP_DIR / job_id
        if job_dir.exists():
            shutil.rmtree(job_dir)
        del jobs[job_id]
    return {"message": "Job cleaned up"}


# Serve static files
if (BASE_DIR / "static").exists():
    app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")


if __name__ == "__main__":
    import uvicorn
    print(f"GEMINI_API_KEY configured: {bool(GEMINI_API_KEY)}")
    uvicorn.run(app, host="0.0.0.0", port=8000)
